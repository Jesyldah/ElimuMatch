"""
Build ElimuMatch Executive Pitch PowerPoint (.pptx)
10 min including live demo | 9 slides
Audience: Impact investors, CSR partners, education foundations
Big idea: We can direct fee support to the students most likely to leave school,
          replacing guesswork with transparent, auditable matching.
Ask: Fund an 8-school Year-1 pilot (~KES 2.0M platform cost)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = r"c:\Users\user1\OneDrive\Documents\MSBA\Capstone"
RF = os.path.join(BASE, "report_figures")
VZ = os.path.join(BASE, "visualizations")

# Portal brand (sponsor portal / Streamlit theme)
INK       = RGBColor(0x14, 0x21, 0x3D)
LEAF      = RGBColor(0x1B, 0x7A, 0x5A)
LEAF_DEEP = RGBColor(0x0F, 0x5C, 0x42)
SUN       = RGBColor(0xF4, 0xB9, 0x42)
SAND      = RGBColor(0xF7, 0xF1, 0xE8)
MIST      = RGBColor(0xE8, 0xEF, 0xE9)
MUTED     = RGBColor(0x5C, 0x6B, 0x73)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
NAVY      = INK
TEAL      = LEAF
GREEN     = LEAF
ORANGE    = SUN
CREAM     = SAND
LGRAY     = MIST
DGRAY     = MUTED

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_para(tf, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p


def img(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kw = {"left": left, "top": top}
        if width:
            kw["width"] = width
        if height:
            kw["height"] = height
        slide.shapes.add_picture(path, **kw)
        return True
    return False


def stripe(slide, label):
    fill_bg(slide, SAND)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.10), LEAF)
    add_rect(slide, Inches(0), Inches(0.10), W, Inches(0.04), SUN)
    add_text(slide, Inches(0.7), Inches(0.22), Inches(12), Inches(0.4),
             label, size=11, color=LEAF, bold=True)


def headline(slide, text):
    add_text(slide, Inches(0.7), Inches(0.7), Inches(12), Inches(0.8),
             text, size=24, bold=True, color=NAVY)


def footnote(slide, text, y=7.0):
    add_text(slide, Inches(0.7), Inches(y), Inches(12), Inches(0.4),
             text, size=10, color=DGRAY)


# ============================================================
# SLIDE 1  TITLE
# Horizontal logic: introduces who we are and what we do
# ============================================================
s = add_blank()
fill_bg(s, INK)
add_rect(s, Inches(0), Inches(0), W, Inches(0.10), LEAF)
add_rect(s, Inches(0), Inches(3.15), W, Inches(0.06), SUN)

add_text(s, Inches(1.5), Inches(1.3), Inches(10), Inches(1.0),
         "ElimuMatch", size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(2.3), Inches(10), Inches(0.7),
         "Matching Support to Retention Risk", size=26,
         color=CREAM, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(3.5), Inches(10), Inches(0.5),
         "A data-driven platform for fee support to Kenyan secondary students",
         size=16, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(4.8), Inches(10), Inches(0.4),
         "Prepared for Impact Investors, CSR Partners, and Education Foundations",
         size=14, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(5.4), Inches(10), Inches(0.4),
         "Jesyldah  |  Founder  |  August 2026",
         size=14, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(6.2), Inches(10), Inches(0.4),
         "jesyldah.github.io/ElimuMatch", size=13,
         color=SUN, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2  THE PROBLEM
# Horizontal logic: Kenya has a retention crisis
# ============================================================
s = add_blank()
fill_bg(s, SAND)
stripe(s, "CONTEXT AND PROBLEM")

headline(s, "Kenya expanded access, but nearly half of students still do not finish secondary school")

img(s, os.path.join(RF, "ext_02_kenya_completion_funnel.png"),
    Inches(0.7), Inches(1.7), width=Inches(5.5))
img(s, os.path.join(RF, "ext_03_kenya_access_progress.png"),
    Inches(6.8), Inches(1.7), width=Inches(5.5))

tb = add_text(s, Inches(0.7), Inches(5.8), Inches(12), Inches(1.0), "", size=15, color=INK)
tf = tb.text_frame
tf.paragraphs[0].text = "4.32 million students were enrolled by 2024, all vulnerable to fee pressure and income shocks."
tf.paragraphs[0].font.size = Pt(15); tf.paragraphs[0].font.color.rgb = INK
add_para(tf, "People want to help, but support follows personal connections instead of reaching those who need it most.", size=15, color=INK)

footnote(s, "Sources: UNESCO IICBA (2025) completion estimates; KNBS Economic Survey 2025, secondary enrolment 2020-2024.")


# ============================================================
# SLIDE 3  THE OPPORTUNITY
# Horizontal logic: current giving misses the mark
# ============================================================
s = add_blank()
fill_bg(s, SAND)
stripe(s, "EVIDENCE OF THE OPPORTUNITY")

headline(s, "NGO help often focuses on ASAL, while students elsewhere also face fee pressure and dropout risk")

img(s, os.path.join(RF, "09_perceptual_map.png"),
    Inches(0.7), Inches(1.7), width=Inches(5.8))

tb = add_text(s, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.5), "", size=15, color=INK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "How support is directed today:"
tf.paragraphs[0].font.size = Pt(16); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = NAVY
add_para(tf, "\u2022  Many NGO programs concentrate on arid and semi-arid lands (ASAL), while students elsewhere also face fee pressure and dropout risk.", size=14, color=INK)
add_para(tf, "\u2022  Most individual donors rely on personal contacts, so students without connections are overlooked.", size=14, color=INK)
add_para(tf, "\u2022  Bursary contests run once a year with heavy paperwork, leaving gaps the rest of the time.", size=14, color=INK)
add_para(tf, "\u2022  Crowdfunding and school lists rarely rank who is most at risk of leaving.", size=14, color=INK)
add_para(tf, "", size=10)
add_para(tf, "ElimuMatch covers all 47 counties: donors can choose any place, and priority follows measured risk, not geography alone.", size=15, color=TEAL, bold=True)

footnote(s, "ASAL focus is a common NGO pattern; ElimuMatch does not claim ASAL work is unnecessary. Bursaries and ElimuMatch are complements.")


# ============================================================
# SLIDE 4  THE SOLUTION + DEMO CUE
# Horizontal logic: here is what we built and how it works
# ============================================================
s = add_blank()
fill_bg(s, SAND)
stripe(s, "OUR SOLUTION")

headline(s, "ElimuMatch identifies at-risk students and lets donors support them in four simple steps")

img(s, os.path.join(RF, "01_product_layers.png"),
    Inches(0.5), Inches(1.6), width=Inches(6.0))

tb = add_text(s, Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.0), "", size=15, color=INK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "What donors experience:"
tf.paragraphs[0].font.size = Pt(16); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = GREEN
add_para(tf, "They pick a county, choose a school, see a student's fee balance, give, and get a receipt.", size=15, color=INK)
add_para(tf, "", size=8)
add_para(tf, "What happens behind the scenes:", size=16, bold=True, color=TEAL)
add_para(tf, "Our system scores each student's risk of dropping out. Only students who need fee help appear. School staff can see the reasons behind each score.", size=15, color=INK)
add_para(tf, "", size=8)
add_para(tf, "Every shilling goes directly to the school's fee account. ElimuMatch does not take a cut.", size=15, color=NAVY, bold=True)

footnote(s, "Privacy: donors see anonymized profiles only. Student identities are protected. Orphan status is never used as an input.")


# ============================================================
# SLIDE 5  HOW IT WORKS (gift journey)
# Horizontal logic: show the donor path visually
# ============================================================
s = add_blank()
fill_bg(s, SAND)
stripe(s, "HOW IT WORKS")

headline(s, "A donor chooses where to give, pays school fees, and gets a receipt")

img(s, os.path.join(RF, "10_gift_journey.png"),
    Inches(1.4), Inches(1.55), height=Inches(3.5))

tb = add_text(s, Inches(0.7), Inches(5.2), Inches(12), Inches(1.2), "", size=15, color=INK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "Every gift goes directly to the school's fee account. Oldest unpaid terms are cleared first."
tf.paragraphs[0].font.size = Pt(15); tf.paragraphs[0].font.color.rgb = INK
add_para(tf, "Overpayments are blocked automatically, and every transaction is tracked and auditable.", size=15, color=INK)

footnote(s, "Live demo: jesyldah.github.io/ElimuMatch")


# ============================================================
# SLIDE 6  WHY IT WORKS
# Horizontal logic: the data backs it up
# ============================================================
s = add_blank()
fill_bg(s, SAND)
stripe(s, "KEY INSIGHT")

headline(s, "Our system identifies about two out of three students who would drop out, twice the nearest alternative")

img(s, os.path.join(RF, "08_selection_rule.png"),
    Inches(0.7), Inches(1.6), width=Inches(7.0))

tb = add_text(s, Inches(8.2), Inches(1.6), Inches(4.5), Inches(4.5), "", size=15, color=INK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "We tested four different approaches."
tf.paragraphs[0].font.size = Pt(15); tf.paragraphs[0].font.color.rgb = INK
add_para(tf, "Logistic Regression is the one we selected: it catches about 67% of students who later drop out. Gradient Boosting catches only about 33%.", size=15, color=INK)
add_para(tf, "", size=10)
add_para(tf, "The strongest warning signs are health problems, family economic pressure, and falling grades.", size=15, color=NAVY, bold=True)
add_para(tf, "", size=10)
add_para(tf, "School staff see why each student is flagged, so they can take informed action. Donors see only a simple, clear profile.", size=14, color=INK)

footnote(s, "Results from proof-of-concept data (1,000 students). Real school validation is the next step. Detection is weaker where nearly all students stay enrolled.")


# ============================================================
# SLIDE 6  BUSINESS VALUE
# Horizontal logic: the numbers work
# ============================================================
s = add_blank()
fill_bg(s, SAND)
stripe(s, "BUSINESS VALUE")

headline(s, "An 8-school pilot returns about KES 2.6 for every KES 1 of platform cost")

img(s, os.path.join(RF, "06_year1_economics.png"),
    Inches(0.5), Inches(1.6), width=Inches(6.0))

tb = add_text(s, Inches(7.0), Inches(1.6), Inches(5.8), Inches(4.5), "", size=15, color=INK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "Year-1 pilot: 8 schools, about 1,000 students"
tf.paragraphs[0].font.size = Pt(17); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = NAVY
add_para(tf, "", size=6)
add_para(tf, "It costs about KES 2.0M to set up and run the platform for one year.", size=15, color=INK)
add_para(tf, "Donor gifts of about KES 4.0M pass through to schools separately.", size=15, color=INK)
add_para(tf, "The estimated benefits total about KES 5.2M from better targeting, fewer dropouts, and staff time saved.", size=15, color=LEAF, bold=True)
add_para(tf, "", size=6)
add_para(tf, "Conservative estimate: roughly break-even.", size=14, color=DGRAY)
add_para(tf, "Base case: 2.6x return.  Best case: 5.4x.", size=16, color=LEAF, bold=True)
add_para(tf, "", size=8)
add_para(tf, "We measure success by whether priority students actually receive gifts, whether payments land without errors, and whether helped students stay in school.", size=14, color=INK)

footnote(s, "All figures are estimates for an illustrative pilot, not audited financials. Benefits depend on real outcomes, not test data.")


# ============================================================
# SLIDE 7  ROADMAP
# Horizontal logic: here is the plan, with built-in discipline
# No image, just phase boxes to avoid overlap
# ============================================================
s = add_blank()
fill_bg(s, SAND)
stripe(s, "IMPLEMENTATION ROADMAP")

headline(s, "We scale only if the pilot clears three decision gates")

phases = [
    ("1. Legal Gate",
     "Sign agreements with\n8 partner schools.\nSet up data protection\nand child safeguarding.",
     "Months 0-2", INK, WHITE),
    ("2. Soft Pilot",
     "Rank students internally.\nReview every list before\nanything is made public.",
     "Term 1", LEAF, WHITE),
    ("3. Live Giving",
     "Donors give through\nthe platform. Money\ngoes directly to\nschool fee accounts.",
     "Terms 1-2", LEAF_DEEP, WHITE),
    ("4. Measure",
     "Track whether helped\nstudents stay in school.\nCheck payment accuracy\nand fairness each term.",
     "Each term", SUN, INK),
    ("5. Scale Decision",
     "Add more schools or\nnew types of support\nonly if the pilot\nresults justify it.",
     "End of Year 1", INK, WHITE),
]
box_w = 2.3; gap = 0.22; x_start = 0.7; y_top = 1.8
for i, (title, desc, timing, color, fg) in enumerate(phases):
    x = Inches(x_start + i * (box_w + gap))
    add_rect(s, x, Inches(y_top), Inches(box_w), Inches(3.45), color)
    add_text(s, x + Inches(0.15), Inches(y_top + 0.15), Inches(box_w - 0.3), Inches(0.4),
             title, size=15, color=fg, bold=True)
    add_text(s, x + Inches(0.15), Inches(y_top + 0.7), Inches(box_w - 0.3), Inches(2.2),
             desc, size=13, color=fg)
    add_text(s, x + Inches(0.15), Inches(y_top + 2.9), Inches(box_w - 0.3), Inches(0.4),
             timing, size=12, color=fg, bold=True)

# Arrows between boxes
for i in range(4):
    x_from = Inches(x_start + i * (box_w + gap) + box_w)
    add_text(s, x_from + Inches(0.02), Inches(y_top + 1.6), Inches(0.2), Inches(0.4),
             "\u25B6", size=14, color=TEAL, align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.7), Inches(5.45), Inches(12.0), Inches(0.42), NAVY)
add_text(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.25),
         "Scale is unlocked only if these three results hold:", size=15, color=WHITE, bold=True)

gate_y = 5.98
gate_w = 3.8
gate_h = 0.62
gate_gap = 0.25
gate_xs = [0.7, 4.75, 8.8]
gate_items = [
    "The model identifies at least 40% of students who later drop out",
    "Payment error rates stay below 10%",
    "Helped students stay in school longer within two terms",
]
for x_val, item in zip(gate_xs, gate_items):
    add_rect(s, Inches(x_val), Inches(gate_y), Inches(gate_w), Inches(gate_h), LGRAY)
    add_text(s, Inches(x_val + 0.1), Inches(gate_y + 0.08), Inches(gate_w - 0.2), Inches(0.45),
             item, size=13, color=INK, bold=True, align=PP_ALIGN.CENTER)

footnote(s, "Every gate requires a formal review. Capital unlocks the next stage only when the previous one clears.")


# ============================================================
# SLIDE 8  THE ASK
# Horizontal logic: here is what we need from you
# ============================================================
s = add_blank()
fill_bg(s, INK)
add_rect(s, Inches(0), Inches(0), W, Inches(0.10), LEAF)
add_rect(s, Inches(0), Inches(2.7), W, Inches(0.06), SUN)

add_text(s, Inches(1.5), Inches(0.9), Inches(10), Inches(0.6),
         "The Ask", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(1.6), Inches(10), Inches(0.9),
         "Fund an 8-school Year-1 pilot to prove\nfee-support matching on real school data",
         size=22, color=CREAM, align=PP_ALIGN.CENTER)

items = [
    "About KES 2.0M for platform setup and operations (donor gifts go to schools separately).",
    "A part-time analyst and school liaison to run the pilot.",
    "Access to partner school records under signed agreements.",
    "Live fee giving with human review of every student list before it is published.",
    "A decision to expand only after pilot results justify it.",
]
tb = add_text(s, Inches(2.5), Inches(3.1), Inches(8), Inches(2.8), "", size=17, color=WHITE)
tf = tb.text_frame; tf.word_wrap = True
for j, item in enumerate(items):
    if j == 0:
        tf.paragraphs[0].text = "\u2713  " + item
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = WHITE
    else:
        add_para(tf, "\u2713  " + item, size=16, color=WHITE, space_before=Pt(10))

add_text(s, Inches(1.0), Inches(6.0), Inches(11), Inches(0.8),
         "We can direct fee support to the students most likely to leave school,\n"
         "replacing guesswork with transparent, auditable matching.",
         size=15, bold=True, color=CREAM, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9  SOURCES
# Horizontal logic: here is our evidence base
# ============================================================
s = add_blank()
fill_bg(s, SAND)
stripe(s, "SOURCES")

add_text(s, Inches(0.7), Inches(0.7), Inches(12), Inches(0.5),
         "References", size=22, bold=True, color=NAVY)

refs = [
    "Adhola, F., Ochola, J., & Tikoko, B. (2025). Donor funding and school operations, Nakuru County.",
    "Adelman, C. (2006). The Toolbox Revisited. U.S. Dept. of Education.",
    "Basch, C. E. (2011). Healthier students are better learners. J. School Health, 81(10).",
    "Childress, M. (2015). Dynamics of education in Kenya. The School Fund.",
    "Cortez, P. & Silva, A. (2008). Using data mining to predict student performance.",
    "Cursor. (2026). Auto (Composer) drafting assistance. cursor.com.",
    "Glennerster, R. et al. (2011). Access and quality in Kenyan education. J-PAL / IPA.",
    "Gongera, E. & Okoth, N. (2013). Alternative financing, Kisii County.",
    "Kenya National Bureau of Statistics. (2025). Economic Survey 2025.",
    "Lundberg, S. & Lee, S.-I. (2017). Interpreting model predictions. NeurIPS 30.",
    "Otieno, M. & Ochieng, J. (2020). 100% transition policy, Machakos.",
    "Realinho, V. et al. (2022). Predicting student dropout and success.",
    "RELI Africa. (2020). Status of secondary education in Kenya.",
    "Tinto, V. (1993). Leaving College (2nd ed.). U. of Chicago Press.",
    "UNESCO IICBA. (2025). Kenya education data brief.",
]
tb = add_text(s, Inches(0.7), Inches(1.3), Inches(5.8), Inches(5.5), "", size=11, color=INK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = refs[0]
tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.color.rgb = INK
for r in refs[1:]:
    add_para(tf, r, size=11, color=INK, space_before=Pt(3))

tb2 = add_text(s, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.5), "", size=11, color=INK)
tf2 = tb2.text_frame; tf2.word_wrap = True
tf2.paragraphs[0].text = "Figure sources:"
tf2.paragraphs[0].font.size = Pt(12); tf2.paragraphs[0].font.bold = True; tf2.paragraphs[0].font.color.rgb = NAVY
add_para(tf2, "Slide 2: Author charts from UNESCO IICBA (2025) and KNBS (2025)", size=11, color=INK)
add_para(tf2, "Slide 3: Illustrative positioning map (author)", size=11, color=INK)
add_para(tf2, "Slide 4: Product layer diagram (author)", size=11, color=INK)
add_para(tf2, "Slide 5: Gift journey and operations flow (author)", size=11, color=INK)
add_para(tf2, "Slide 6: Model comparison from project pipeline on test data", size=11, color=INK)
add_para(tf2, "Slide 7: Year-1 economics chart (illustrative, author)", size=11, color=INK)
add_para(tf2, "", size=8)
add_para(tf2, "AI-assisted drafting:", size=12, bold=True, color=NAVY)
add_para(tf2, "Report drafting supported by Cursor Auto (Composer) (Cursor, 2026). All decisions and claims are the author's.", size=11, color=INK)
add_para(tf2, "", size=8)
add_para(tf2, "All results are from proof-of-concept data (1,000 students). Not yet validated on live school records.", size=11, color=DGRAY, bold=True)

add_text(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
         "Live demo: jesyldah.github.io/ElimuMatch  |  Code: github.com/Jesyldah/ElimuMatch",
         size=11, color=TEAL)


# ============================================================
out = os.path.join(BASE, "ElimuMatch_Executive_Pitch.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")

# Print horizontal logic check
titles = [
    "ElimuMatch: A data-driven platform for fee support to Kenyan secondary students",
    "Kenya expanded access, but nearly half of students still do not finish secondary school",
    "NGO help often focuses on ASAL, while students elsewhere also face fee pressure and dropout risk",
    "ElimuMatch identifies at-risk students and lets donors support them in four simple steps",
    "A donor chooses where to give, pays school fees, and gets a receipt",
    "Our system identifies about two out of three students who would drop out, twice the nearest alternative",
    "An 8-school pilot returns about KES 2.6 for every KES 1 of platform cost",
    "We scale only if the pilot clears three decision gates",
    "Fund an 8-school Year-1 pilot to prove fee-support matching on real school data",
    "Sources",
]
print("\nHorizontal logic (read headlines only):")
for i, t in enumerate(titles, 1):
    print(f"  {i}. {t}")
